from typing import Optional


class PasswordProtectedError(Exception):
    """Raised when a file cannot be read because it is password-protected."""


SUPPORTED_IMAGE_TYPES = [
    "image/jpeg",
    "image/png",
    "image/gif",
    "image/webp",
]

SUPPORTED_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "txt": "text/plain",
}


def extract_text_from_file(file_path: str, mime_type: str) -> dict:
    """
    Extract text from a file based on its MIME type.

    Args:
        file_path: Path to the file
        mime_type: MIME type of the file

    Returns:
        dict with success status, text, and metadata
    """
    try:
        text = None
        file_type = None

        if mime_type == "application/pdf" or file_path.endswith(".pdf"):
            text = _extract_from_pdf(file_path)
            file_type = "pdf"

        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or file_path.endswith(".docx")
        ):
            text = _extract_from_docx(file_path)
            file_type = "docx"

        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or file_path.endswith(".pptx")
        ):
            text = _extract_from_pptx(file_path)
            file_type = "pptx"

        elif mime_type == "text/plain" or file_path.endswith(".txt"):
            text = _extract_from_txt(file_path)
            file_type = "txt"

        else:
            return {
                "success": False,
                "error": f"Unsupported file type: {mime_type}",
                "can_vectorize": False,
            }

        if text is None or len(text.strip()) < 10:
            return {
                "success": False,
                "error": "Could not extract text from file or file is empty",
                "can_vectorize": False,
            }

        return {
            "success": True,
            "text": text,
            "file_type": file_type,
            "char_count": len(text),
            "can_vectorize": True,
        }

    except PasswordProtectedError:
        return {
            "success": False,
            "error": "password_protected",
            "password_protected": True,
            "can_vectorize": False,
        }

    except Exception as e:
        print(f"Error extracting text from file: {e}")
        return {
            "success": False,
            "error": str(e),
            "can_vectorize": False,
        }


def _extract_from_pdf(file_path: str) -> Optional[str]:
    """
    Extract text from PDF using PyPDF2.
    Each page is prefixed with <<PAGE:N>> so the chunker can track page ranges.
    """
    import PyPDF2

    parts = []
    with open(file_path, "rb") as f:
        pdf_reader = PyPDF2.PdfReader(f)
        if pdf_reader.is_encrypted:
            raise PasswordProtectedError("PDF is password protected")
        for i, page in enumerate(pdf_reader.pages, start=1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                parts.append(f"<<PAGE:{i}>>\n{page_text}")
    text = "\n".join(parts)
    return text if text.strip() else None


def _extract_from_docx(file_path: str) -> Optional[str]:
    """Extract text from DOCX using python-docx."""
    import zipfile
    from docx import Document

    try:
        doc = Document(file_path)
    except zipfile.BadZipFile:
        raise PasswordProtectedError("DOCX is password protected")
    text = "\n".join([para.text for para in doc.paragraphs])
    return text if text.strip() else None


def _extract_from_pptx(file_path: str) -> Optional[str]:
    """Extract text from PPTX using python-pptx."""
    import zipfile
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
    except zipfile.BadZipFile:
        raise PasswordProtectedError("PPTX is password protected")
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
    return text if text.strip() else None


def _extract_from_txt(file_path: str) -> Optional[str]:
    """Extract text from plain text file."""
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()
    return text if text.strip() else None


def get_file_extension(mime_type: str) -> str:
    """Get file extension from MIME type."""
    extension_map = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "text/plain": "txt",
        "image/jpeg": "jpg",
        "image/png": "png",
        "image/gif": "gif",
        "image/webp": "webp",
    }
    return extension_map.get(mime_type, "bin")


def is_vectorizable(mime_type: str) -> bool:
    """Check if a file type can be vectorized."""
    vectorizable_types = [
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "text/plain",
    ]
    return mime_type in vectorizable_types


def extract_image_description(file_path: str, mime_type: str) -> dict:
    """
    Extract a text description from an image using Claude's vision API.
    Describes visual content and extracts any visible text for RAG indexing.

    Returns:
        dict with success status, text description, and can_vectorize flag
    """
    import base64
    import anthropic
    from agente_rolplay.config import ANTHROPIC_API_KEY, ANTHROPIC_MODEL_NAME, IMAGE_DESCRIPTION_MAX_TOKENS

    if mime_type not in SUPPORTED_IMAGE_TYPES:
        return {"success": False, "error": f"Unsupported image type: {mime_type}", "can_vectorize": False}

    try:
        with open(file_path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode("utf-8")

        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

        response = client.messages.create(
            model=ANTHROPIC_MODEL_NAME,
            max_tokens=IMAGE_DESCRIPTION_MAX_TOKENS,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": (
                                "Describe this image in detail for a searchable knowledge base. "
                                "Include: all visible text (OCR), visual elements, charts, diagrams, "
                                "tables, people, objects, colors, and any other relevant details. "
                                "Be comprehensive so the image can be found via semantic text search."
                            ),
                        },
                    ],
                }
            ],
        )

        description = response.content[0].text

        if not description or len(description.strip()) < 10:
            return {"success": False, "error": "No description generated", "can_vectorize": False}

        return {
            "success": True,
            "text": description,
            "file_type": "image",
            "char_count": len(description),
            "can_vectorize": True,
        }

    except Exception as e:
        print(f"Error extracting image description via vision: {e}")
        return {"success": False, "error": str(e), "can_vectorize": False}


def get_file_type_category(mime_type: str) -> str:
    """Get the category of a file type."""
    if mime_type.startswith("image/"):
        return "image"
    elif mime_type in [
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "text/plain",
    ]:
        return "document"
    else:
        return "other"
