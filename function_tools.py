# python3 function_tools.py

from anthropic import Anthropic
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from qdrant_client import models, QdrantClient
from qdrant_client.models import PointStruct
from system_prompt import prompt_clasificador_saludo_inicial

import json
import os
# import pandas as pd
import pytz
import time

load_dotenv()
anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
MODEL_NAME = os.getenv("ANTHROPIC_MODEL_NAME")
openai_api_key = os.getenv('OPENAI_API_KEY')
OPENAI_EMBEDDINGS_MODEL = os.getenv('OPENAI_EMBEDDINGS_MODEL')
VECTOR_DIMENSION = int(os.getenv('VECTOR_DIMENSION', 1024))
N_SIMILARITY = int(os.getenv('N_SIMILARITY', 3))
QDRANT_URL = os.getenv("QDRANT_URL")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY")
QDRANT_COLLECTION_NAME = os.getenv("QDRANT_COLLECTION_NAME")

client = Anthropic(api_key=anthropic_api_key)
openai_client = OpenAI(api_key=openai_api_key)
qdrant_client = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)

dict_clientes_datos = {}

def anthropic_completion(system_prompt, messages, model_name=MODEL_NAME, max_tokens=5192, temperature=0.0, client=client):
    answer = client.messages.create(
        model=model_name,
        messages=messages,
        max_tokens=max_tokens,
        system=system_prompt,
        temperature=temperature
    )
    return answer

def create_embeddings(texto, client=openai_client, model=OPENAI_EMBEDDINGS_MODEL, dim=VECTOR_DIMENSION):
    respuesta = client.embeddings.create(
        input=texto,
        model=model,
        dimensions=dim,
    )
    return {
        'answer':respuesta.data[0].embedding
    }

def insert_info_business(secciones, client_qdrant=qdrant_client, COLECCION = QDRANT_COLLECTION_NAME, dim=VECTOR_DIMENSION):

    collections = client_qdrant.get_collections().collections
    collection_names = [collection.name for collection in collections]

    if COLECCION not in collection_names:
        client_qdrant.create_collection(
            collection_name=COLECCION,
            vectors_config={
                "embeddings": {
                    "size": dim,
                    "distance": "Cosine"
                }
            }
        )
        print(f"Colección '{COLECCION}' creada exitosamente.")

    puntos = []

    for index, seccion in enumerate(secciones):

        vector = create_embeddings(seccion["texto"]) 

        punto = PointStruct(
            id=index,
            vector={"embeddings": vector['answer']},
            payload={
                "nombre": seccion["nombre"],
                "texto": seccion["texto"]
            }
        )

        puntos.append(punto)

    client_qdrant.upsert(collection_name=COLECCION, points=puntos)
    print(f"{len(puntos)} secciones insertadas en Qdrant.")

def get_text_by_relevance(consulta, cliente=qdrant_client, coleccion=QDRANT_COLLECTION_NAME, n=N_SIMILARITY):
    """
    Busca textos relevantes en Qdrant usando query_points (nueva API)
    """
    try:
        # Generar embedding de la consulta
        embedding = create_embeddings(consulta)
        
        # 🔥 NUEVA API: query_points en lugar de search
        resultado_busqueda = cliente.query_points(
            collection_name=coleccion,
            query=embedding['answer'],  # Vector directo, sin tupla
            using="embeddings",  # Nombre del vector
            limit=n,
            with_payload=True
        )
        
        print(f"✅ Búsqueda exitosa: {len(resultado_busqueda.points)} resultados")
        
        # Procesar resultados
        texto_relevante = []
        for scored_point in resultado_busqueda.points:
            texto_relevante.append({
                'texto': scored_point.payload.get('texto', ''),
                'nombre': scored_point.payload.get('nombre', ''),
                'ruta': scored_point.payload.get('ruta', ''),
                'file_id': scored_point.payload.get('file_id', ''),
                'score': scored_point.score
            })
        
        return texto_relevante
        
    except Exception as e:
        print(f"❌ Error en búsqueda de Qdrant: {e}")
        import traceback
        traceback.print_exc()
        return []

def get_mexico_city_time():
    timestamp = time.time()
    utc_time = datetime.utcfromtimestamp(timestamp)
    mexico_city_timezone = pytz.timezone('America/Mexico_City')
    mexico_city_time = utc_time.astimezone(mexico_city_timezone)
    formatted_time = mexico_city_time.strftime('%A, %Y-%m-%d %H:%M')
    return formatted_time

def agregar_punto_individual(texto, nombre, client_qdrant=qdrant_client, COLECCION=QDRANT_COLLECTION_NAME):
    try:
        # Obtener el último ID de la colección
        resultado = client_qdrant.scroll(
            collection_name=COLECCION,
            limit=1,
            with_payload=False,
            with_vectors=False,
            order_by="id"
        )

        # Si hay puntos, tomar el último ID y sumar 1, sino empezar en 0
        if resultado[0]:
            ultimo_id = max([punto.id for punto in resultado[0]]) if resultado[0] else -1
            nuevo_id = ultimo_id + 1
        else:
            nuevo_id = 0

        vector = create_embeddings(texto)

        punto = PointStruct(
            id=nuevo_id,
            vector={"embeddings": vector['answer']},
            payload={
                "nombre": nombre,
                "texto": texto
            }
        )

        client_qdrant.upsert(collection_name=COLECCION, points=[punto])

        return {"success": True, "id": nuevo_id, "message": "Información agregada de forma exitosa."}

    except Exception as e:
        return {"success": False, "message": f"Error: {str(e)}"}

def insert_datos_pauta(secciones, client_qdrant=qdrant_client, COLECCION=QDRANT_COLLECTION_NAME, dim=VECTOR_DIMENSION):
    collections = client_qdrant.get_collections().collections
    collection_names = [collection.name for collection in collections]

    if COLECCION not in collection_names:
        client_qdrant.create_collection(
            collection_name=COLECCION,
            vectors_config={
                "embeddings": {
                    "size": dim,
                    "distance": "Cosine"
                }
            }
        )
        print(f"Colección '{COLECCION}' creada exitosamente.")

    puntos = []
    for index, seccion in enumerate(secciones):
        vector = create_embeddings(seccion["texto_embeddings"]) 

        punto = PointStruct(
            id=index,
            vector={"embeddings": vector['answer']},
            payload=seccion['datos_completos_punto']
            # payload={
            #     "nombre": seccion["nombre"],
            #     "texto": seccion["texto"]
            # }
        )
        puntos.append(punto)

    client_qdrant.upsert(collection_name=COLECCION, points=puntos)
    print(f"{len(puntos)} secciones insertadas en Qdrant.")

# ======================================
def insertar_documentos_drive_a_qdrant(
    json_path="textos_drive_extraidos.json",
    client_qdrant=qdrant_client, 
    COLECCION=QDRANT_COLLECTION_NAME, 
    dim=VECTOR_DIMENSION
):
    """
    Inserta los documentos extraídos de Google Drive en Qdrant
    
    Args:
        json_path (str): Ruta al archivo JSON con los textos extraídos
        client_qdrant: Cliente de Qdrant
        COLECCION (str): Nombre de la colección
        dim (int): Dimensión de los embeddings
    """
    collections = client_qdrant.get_collections().collections
    collection_names = [collection.name for collection in collections]

    if COLECCION not in collection_names:
        client_qdrant.create_collection(
            collection_name=COLECCION,
            vectors_config={
                "embeddings": {
                    "size": dim,
                    "distance": "Cosine"
                }
            }
        )
        print(f"✅ Colección '{COLECCION}' creada exitosamente.")

    else:
        print(f"ℹ️ Colección '{COLECCION}' ya existe.")

    print(f"📂 Cargando documentos desde: {json_path}")
    with open(json_path, 'r', encoding='utf-8') as f:
        documentos = json.load(f)
    
    print(f"📚 Total de documentos a procesar: {len(documentos)}")
    
    puntos = []
    
    for index, doc in enumerate(documentos):
        print(f"🔄 Procesando [{index + 1}/{len(documentos)}]: {doc['nombre']}")
        
        # Preparar el texto para embeddings (nombre + texto)
        texto_para_embedding = f"{doc['nombre']}\n\n{doc['texto']}"
        
        # Limitar a los primeros 8000 caracteres para no exceder límites de tokens
        texto_truncado = texto_para_embedding[:8000]
        
        # Generar embedding
        try:
            vector = create_embeddings(texto_truncado)
            
            # Crear punto con payload completo
            punto = PointStruct(
                id=index,
                vector={"embeddings": vector['answer']},
                payload={
                    "file_id": doc['file_id'],
                    "nombre": doc['nombre'],
                    "ruta": doc['ruta'],
                    "mime_type": doc['mime_type'],
                    "texto": doc['texto'],  # Texto completo en payload
                    "num_caracteres": doc['num_caracteres'],
                    "num_palabras": doc['num_palabras']
                }
            )
            
            puntos.append(punto)
            print(f"  ✅ Embedding generado para: {doc['nombre']}")
            
        except Exception as e:
            print(f"  ❌ Error generando embedding para {doc['nombre']}: {e}")
            continue
    
    # Insertar todos los puntos en Qdrant
    if puntos:
        print(f"\n📤 Insertando {len(puntos)} documentos en Qdrant...")
        client_qdrant.upsert(collection_name=COLECCION, points=puntos)
        print(f"✅ {len(puntos)} documentos insertados exitosamente en Qdrant.")
    else:
        print("⚠️ No hay puntos para insertar.")
    
    return len(puntos)

def agregar_documento_a_qdrant(file_id, mime_type, nombre, ruta, ruta_temporal):
    """
    Extrae texto y agrega documento a Qdrant
    Soporta: PDF, Google Docs, DOCX
    """
    try:
        print(f"📄 Extrayendo texto de: {nombre}")
        texto = None
        
        # PDF
        if mime_type == 'application/pdf':
            import PyPDF2
            with open(ruta_temporal, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                texto = ""
                for page in pdf_reader.pages:
                    texto += page.extract_text() + "\n"
        
        # DOCX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            from docx import Document
            doc = Document(ruta_temporal)
            texto = "\n".join([para.text for para in doc.paragraphs])
        
        # 🔥 PPTX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            from pptx import Presentation
            prs = Presentation(ruta_temporal)
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"

        # 🔥 XLSX
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            import openpyxl
            wb = openpyxl.load_workbook(ruta_temporal)
            texto = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                texto += f"\n=== Hoja: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    texto += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"

        
        # Google Doc (exportar desde Drive)
        elif mime_type == 'application/vnd.google-apps.document':
            from google_drive import obtener_servicio
            service = obtener_servicio()
            request = service.files().export_media(fileId=file_id, mimeType='text/plain')
            from io import BytesIO
            from googleapiclient.http import MediaIoBaseDownload
            fh = BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
            fh.seek(0)
            texto = fh.read().decode('utf-8')
        
        # SI NO ES NINGUNO DE ESTOS
        else:
            print(f"⚠️ Tipo de archivo no soportado para extracción: {mime_type}")
            return False
        
        print(F" TEXTO: ", texto)
        if not texto or len(texto.strip()) < 10:
            print(f"⚠️ Texto muy corto o vacío")
            return False
        
        print(f"✅ Texto extraído: {len(texto)} caracteres")
        
        # Generar embedding
        texto_para_embedding = f"{nombre}\n\n{texto[:8000]}"
        embedding = create_embeddings(texto_para_embedding)
        
        # Crear punto
        punto_id = int(time.time())
        punto = PointStruct(
            id=punto_id,
            vector={"embeddings": embedding['answer']},
            payload={
                "file_id": file_id,
                "nombre": nombre,
                "ruta": ruta,
                "mime_type": mime_type,
                "texto": texto,
                "num_caracteres": len(texto),
                "fecha_subida": datetime.now(pytz.timezone('America/Mexico_City')).isoformat()
            }
        )
        
        # Insertar en Qdrant
        qdrant_client.upsert(collection_name=QDRANT_COLLECTION_NAME, points=[punto])
        print(f"✅ Agregado a Qdrant con ID: {punto_id}")
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return False

if __name__ == "__main__":

    insertar_documentos_drive_a_qdrant()

    # Proceso para insertar los datos del spreadsheets en Qdrant
    # file_path = "datos_pauta.xlsx"
    # datos_embeddear = read_spreadsheets_data_and_generate_dict_embeds(file_path)
    # insert_datos_pauta(datos_embeddear)

    # Proceso para probar el obtener el texto por relevancia
    # consulta = "dime sobre los clientes de san fer ? "
    # consulta = "dame informacion sobre las minutas que hubo el 16 de octubre del 25   "
    # ans = get_text_by_relevance(consulta)
    # print("ans", ans)