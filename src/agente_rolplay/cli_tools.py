# python3 cli_tools.py (formerly function_tools.py)

from anthropic import Anthropic
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from pinecone import Pinecone, ServerlessSpec
from src.agente_rolplay.system_prompt import (
    prompt_clasificador_saludo_inicial,
    system_prompt_categorize,
)

import json
import os

# import pandas as pd
import pytz
import time

load_dotenv()
anthropic_api_key = os.getenv("ANTHROPIC_API_KEY")
MODEL_NAME = os.getenv("ANTHROPIC_MODEL_NAME")
openai_api_key = os.getenv("OPENAI_API_KEY")
OPENAI_EMBEDDINGS_MODEL = os.getenv("OPENAI_EMBEDDINGS_MODEL")
VECTOR_DIMENSION = int(os.getenv("VECTOR_DIMENSION", 1024))
N_SIMILARITY = int(os.getenv("N_SIMILARITY", 3))
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENVIRONMENT = os.getenv("PINECONE_ENVIRONMENT", "us-east-1")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME", "rolplay-knowledge")

client = Anthropic(api_key=anthropic_api_key)
openai_client = OpenAI(api_key=openai_api_key)
pinecone_client = Pinecone(api_key=PINECONE_API_KEY)


def get_pinecone_index():
    """Get or create Pinecone index"""
    try:
        if PINECONE_INDEX_NAME not in pinecone_client.list_indexes().names():
            pinecone_client.create_index(
                name=PINECONE_INDEX_NAME,
                dimension=VECTOR_DIMENSION,
                metric="cosine",
                spec=ServerlessSpec(cloud="aws", region=PINECONE_ENVIRONMENT),
            )
        return pinecone_client.Index(PINECONE_INDEX_NAME)
    except Exception as e:
        print(f"Error with Pinecone index: {e}")
        return None


index = get_pinecone_index()

dict_clientes_datos = {}


def anthropic_completion(
    system_prompt,
    messages,
    model_name=MODEL_NAME,
    max_tokens=5192,
    temperature=0.0,
    client=client,
):
    answer = client.messages.create(
        model=model_name,
        messages=messages,
        max_tokens=max_tokens,
        system=system_prompt,
        temperature=temperature,
    )
    return answer


def create_embeddings(
    text, client=openai_client, model=OPENAI_EMBEDDINGS_MODEL, dim=VECTOR_DIMENSION
):
    response = client.embeddings.create(
        input=text,
        model=model,
        dimensions=dim,
    )
    return {"answer": response.data[0].embedding}


def categorize_document(text, max_chars=3000):
    """
    Categorize a document using the LLM based on its content.
    Returns the category as a string (e.g., 'proposal', 'service', 'meeting').
    """
    try:
        truncated_text = text[:max_chars]

        response = client.messages.create(
            model=MODEL_NAME,
            max_tokens=50,
            system=system_prompt_categorize,
            messages=[
                {
                    "role": "user",
                    "content": f"Clasifica este documento:\n\n{truncated_text}",
                }
            ],
            temperature=0.0,
        )

        category = response.content[0].text.strip().lower()

        valid_categories = [
            "proposal",
            "service",
            "integration",
            "meeting",
            "contract",
            "invoice",
            "technical",
            "company_info",
            "contact",
            "faq",
            "policy",
            "project",
            "other",
        ]

        if category in valid_categories:
            return category
        else:
            return "other"

    except Exception as e:
        print(f"Error categorizing document: {e}")
        return "other"


def insert_info_business(
    sections,
    index_obj=index,
    dim=VECTOR_DIMENSION,
):
    """Insert business information into Pinecone"""
    if index_obj is None:
        print("Pinecone index not available")
        return

    vectors = []
    for index, section in enumerate(sections):
        vector = create_embeddings(section["texto"])
        vectors.append(
            {
                "id": str(index),
                "values": vector["answer"],
                "metadata": {"nombre": section["nombre"], "texto": section["texto"]},
            }
        )

    index_obj.upsert(vectors=vectors)
    print(f"{len(vectors)} sections inserted into Pinecone.")


def get_text_by_relevance(query, index_obj=index, n=N_SIMILARITY):
    """
    Search relevant texts in Pinecone using vector similarity
    """
    if index_obj is None:
        print("Pinecone index not available")
        return []

    try:
        embedding = create_embeddings(query)

        search_result = index_obj.query(
            vector=embedding["answer"],
            top_k=n,
            include_metadata=True,
        )

        print(f"Search successful: {len(search_result['matches'])} results")

        relevant_texts = []
        for match in search_result["matches"]:
            relevant_texts.append(
                {
                    "texto": match["metadata"].get("texto", ""),
                    "nombre": match["metadata"].get("nombre", ""),
                    "ruta": match["metadata"].get("ruta", ""),
                    "file_id": match["metadata"].get("file_id", ""),
                    "score": match["score"],
                    "category": match["metadata"].get("category", "other"),
                }
            )

        return relevant_texts

    except Exception as e:
        print(f"Error in Pinecone search: {e}")
        import traceback

        traceback.print_exc()
        return []


def get_mexico_city_time():
    timestamp = time.time()
    utc_time = datetime.utcfromtimestamp(timestamp)
    mexico_city_timezone = pytz.timezone("America/Mexico_City")
    mexico_city_time = utc_time.astimezone(mexico_city_timezone)
    formatted_time = mexico_city_time.strftime("%A, %Y-%m-%d %H:%M")
    return formatted_time


def agregar_punto_individual(texto, nombre, index_obj=index):
    """Add a single point to Pinecone"""
    try:
        if index_obj is None:
            return {"success": False, "message": "Pinecone index not available"}

        import uuid

        new_id = str(uuid.uuid4())
        vector = create_embeddings(texto)

        index_obj.upsert(
            vectors=[
                {
                    "id": new_id,
                    "values": vector["answer"],
                    "metadata": {"nombre": nombre, "texto": texto},
                }
            ]
        )

        return {
            "success": True,
            "id": new_id,
            "message": "Information added successfully.",
        }

    except Exception as e:
        return {"success": False, "message": f"Error: {str(e)}"}


def insert_datos_pauta(sections, index_obj=index, dim=VECTOR_DIMENSION):
    """Insert pauta data into Pinecone"""
    if index_obj is None:
        print("Pinecone index not available")
        return

    vectors = []
    for index, section in enumerate(sections):
        vector = create_embeddings(section["texto_embeddings"])
        vectors.append(
            {
                "id": str(index),
                "values": vector["answer"],
                "metadata": section["datos_completos_punto"],
            }
        )

    index_obj.upsert(vectors=vectors)
    print(f"{len(vectors)} sections inserted into Pinecone.")


def insertar_documentos_drive_a_qdrant(
    json_path="textos_drive_extraidos.json",
    index_obj=index,
    dim=VECTOR_DIMENSION,
):
    """
    Insert documents extracted from Google Drive into Pinecone

    Args:
        json_path (str): Path to JSON file with extracted texts
        index_obj: Pinecone index
        dim (int): Embedding dimension
    """
    if index_obj is None:
        print("Pinecone index not available")
        return 0

    print(f"Loading documents from: {json_path}")
    with open(json_path, "r", encoding="utf-8") as f:
        documents = json.load(f)

    print(f"Total documents to process: {len(documents)}")

    vectors = []

    for index, doc in enumerate(documents):
        print(f"Processing [{index + 1}/{len(documents)}]: {doc['nombre']}")

        # Prepare text for embeddings (name + text)
        text_for_embedding = f"{doc['nombre']}\n\n{doc['texto']}"

        # Limit to first 8000 characters to not exceed token limits
        truncated_text = text_for_embedding[:8000]

        # Generate embedding
        try:
            vector = create_embeddings(truncated_text)

            # Auto-categorize the document
            print(f"  Categorizing document...")
            category = categorize_document(doc["texto"])
            print(f"  Document category: {category}")

            # Create vector with metadata
            vectors.append(
                {
                    "id": str(index),
                    "values": vector["answer"],
                    "metadata": {
                        "file_id": doc["file_id"],
                        "nombre": doc["nombre"],
                        "ruta": doc["ruta"],
                        "mime_type": doc["mime_type"],
                        "texto": doc["texto"],
                        "num_caracteres": doc["num_caracteres"],
                        "num_palabras": doc["num_palabras"],
                        "category": category,
                    },
                }
            )
            print(f"  Embedding generated for: {doc['nombre']}")

        except Exception as e:
            print(f"  Error generating embedding for {doc['nombre']}: {e}")
            continue

    # Insert all vectors into Pinecone
    if vectors:
        print(f"\nInserting {len(vectors)} documents into Pinecone...")
        index_obj.upsert(vectors=vectors)
        print(f"{len(vectors)} documents inserted successfully into Pinecone.")
    else:
        print("No vectors to insert.")

    return len(vectors)


def agregar_documento_a_qdrant(file_id, mime_type, nombre, ruta, ruta_temporal):
    """
    Extract text and add document to Qdrant
    Supports: PDF, Google Docs, DOCX
    """
    try:
        print(f"Extracting text from: {nombre}")
        text = None

        # PDF
        if mime_type == "application/pdf":
            import PyPDF2

            with open(ruta_temporal, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"

        # DOCX
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            from docx import Document

            doc = Document(ruta_temporal)
            text = "\n".join([para.text for para in doc.paragraphs])

        # PPTX
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            from pptx import Presentation

            prs = Presentation(ruta_temporal)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        # XLSX
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            import openpyxl

            wb = openpyxl.load_workbook(ruta_temporal)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    text += (
                        " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                    )

        # Google Doc - Not supported without Google Drive
        # Users should convert Google Docs to PDF/DOCX before uploading
        elif mime_type == "application/vnd.google-apps.document":
            print("Google Docs require conversion to PDF/DOCX before uploading")
            return False

        # IF NONE OF THESE
        else:
            print(f"Unsupported file type for extraction: {mime_type}")
            return False

        print(f" TEXT: ", text)
        if not text or len(text.strip()) < 10:
            print(f"Text too short or empty")
            return False

        print(f"Text extracted: {len(text)} characters")

        # Auto-categorize the document
        print("Categorizing document...")
        category = categorize_document(text)
        print(f"Document category: {category}")

        # Generate embedding
        text_for_embedding = f"{nombre}\n\n{text[:8000]}"
        embedding = create_embeddings(text_for_embedding)

        # Create vector with metadata
        import uuid

        vector_id = str(uuid.uuid4())

        # Insert into Pinecone
        index.upsert(
            vectors=[
                {
                    "id": vector_id,
                    "values": embedding["answer"],
                    "metadata": {
                        "file_id": file_id,
                        "nombre": nombre,
                        "ruta": ruta,
                        "mime_type": mime_type,
                        "texto": text,
                        "num_caracteres": len(text),
                        "category": category,
                        "fecha_subida": datetime.now(
                            pytz.timezone("America/Mexico_City")
                        ).isoformat(),
                    },
                }
            ]
        )
        print(f"Added to Pinecone with ID: {vector_id}")
        return True

    except Exception as e:
        print(f"Error: {e}")
        return False


if __name__ == "__main__":
    insertar_documentos_drive_a_qdrant()

    # Process to insert spreadsheet data into Pinecone
    # file_path = "datos_pauta.xlsx"
    # datos_embeddear = read_spreadsheets_data_and_generate_dict_embeds(file_path)
    # insert_datos_pauta(datos_embeddear)

    # Process to test getting text by relevance
    # consulta = "dime sobre los clientes de san fer ? "
    # consulta = "dame informacion sobre las minutas que hubo el 16 de octubre del 25   "
    # ans = get_text_by_relevance(consulta)
    # print("ans", ans)
